import io
import os
import re
import json
import uuid
import base64
import sqlite3
import hashlib
import datetime
import copy
import time
from urllib.parse import urljoin, urlparse, parse_qs, unquote, quote
from contextlib import contextmanager
from typing import List, Dict, Any, Optional
from html import escape as html_escape

import requests
import yfinance as yf
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import speech_recognition as sr
import streamlit as st
import streamlit.components.v1 as components
import extra_streamlit_components as stx
from openai import OpenAI
from audio_recorder_streamlit import audio_recorder

# ==========================================
# 1. KONSTANTA & PENGATURAN AI AGENT
# ==========================================
DB_NAME = 'lagos_multiuser.db'
API_KEY = st.secrets["NVIDIA_API_KEY"]
BASE_URL = "https://integrate.api.nvidia.com/v1"

B3 = "`" * 3

MODEL_MAPPING = {
    "meta/muse-glimmer-30b": "1. Aether (Flash)",
    "google/diffusiongemma-26b-a4b-it": "2. Verper (pro)",
    "nvidia/nemotron-3.5-lightning-30b-a3b": "3. Numayr(Eksklusif)",
    "thinkingmachines/inkling": "4. Nova (Unstable)",
    "deepseek-ai/deepseek-v4-flash-0731": "5. Zeta (Under Construction)",
    "google/veo-3.1-fast-generate-preview": "6. Generator Gambar (coming soon)"
}

HTTP = requests.Session()
HTTP.headers.update({
    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36',
    'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8',
    'Accept-Language': 'id-ID,id;q=0.9,en-US;q=0.8,en;q=0.7',
})

def panggil_api_dengan_retry(client_instance, **kwargs):
    max_retries = 4
    for attempt in range(max_retries):
        try:
            return client_instance.chat.completions.create(**kwargs)
        except Exception as e:
            error_msg = str(e)
            if "429" in error_msg and attempt < max_retries - 1:
                jeda = 3 + (attempt * 3)
                st.toast(f"⏳ Menyesuaikan limit 40 RPM. Melanjutkan dalam {jeda} detik... ({attempt+1}/{max_retries})")
                time.sleep(jeda)
            else:
                raise e

LAGOS_TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "ambil_data_pasar",
            "description": "Ambil data harga saham (.JK) atau kripto (-USD).",
            "parameters": {
                "type": "object",
                "properties": {
                    "simbol_ticker": {"type": "string", "description": "Simbol saham (.JK) atau kripto (-USD)."}
                },
                "required": ["simbol_ticker"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "cari_informasi_web",
            "description": "Cari berita/fakta/informasi terkini dari internet. Gunakan kata kunci yang sangat spesifik (termasuk nama daerah jika itu tempat lokal).",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Kata kunci spesifik & akurat."}
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "baca_isi_website",
            "description": "Baca isi lengkap halaman dari URL dalam bentuk teks bersih.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "URL website (dimulai http/https)."}
                },
                "required": ["url"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "cari_gambar",
            "description": "Cari URL foto/gambar asli dari suatu benda, tempat, hewan, atau tokoh.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Nama entitas yang ingin dicari fotonya."}
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "ambil_transkrip_youtube",
            "description": "Ambil teks/transkrip dari URL video YouTube untuk dirangkum.",
            "parameters": {
                "type": "object",
                "properties": {
                    "video_url": {"type": "string", "description": "Link URL video YouTube lengkap."}
                },
                "required": ["video_url"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "eksekusi_python",
            "description": "Jalankan skrip Python murni (analisis data, logika, kalkulasi).",
            "parameters": {
                "type": "object",
                "properties": {
                    "kode": {"type": "string", "description": "Kode Python murni yang ingin dieksekusi."}
                },
                "required": ["kode"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "hitung_matematika",
            "description": "Hitung operasi matematika agar hasilnya akurat.",
            "parameters": {
                "type": "object",
                "properties": {
                    "ekspresi": {"type": "string", "description": "Ekspresi matematika (contoh: 25000 * 1.15)."}
                },
                "required": ["ekspresi"]
            }
        }
    }
]

SYSTEM_PROMPT = """Anda adalah Lagøs AI, Asisten Kecerdasan Buatan analitik tingkat tinggi yang dikembangkan oleh Rian Dev.

KARAKTERISTIK UTAMA ANDA:
1. Anda sangat cerdas, natural, komprehensif, dan suportif layaknya asisten AI kelas dunia (berperilaku mirip Gemini).
2. Anda WAJIB merangkum dan menyintesis informasi dari alat (tools) ke dalam kalimat yang enak dibaca. JANGAN PERNAH memberikan hasil log mentah atau memuntahkan teks mesin pencari ke pengguna. Gunakan gaya bahasa Anda sendiri.
3. Jawab langsung ke inti pertanyaan. Jika diminta menganalisis, berikan analisis yang logis dan terstruktur.
4. Jangan halusinasi. Jika informasi benar-benar tidak ada di internet, katakan jujur dengan sopan bahwa informasinya tidak ditemukan, lalu berikan saran alternatif.
5. Gunakan format Markdown (Bold, Bullet Points) agar teks mudah dibaca. Sertakan tautan sumber (URL) jika mengutip berita/artikel.

ATURAN MULTIMEDIA & FORMAT:
- FOTO ASLI: `![Deskripsi](URL_Gambar)`
- GAMBAR BUATAN AI: `![Generate Gambar](https://image.pollinations.ai/prompt/deskripsi_gambar_dalam_bahasa_inggris_detail?width=800&height=600&nologo=true)`
- PPT JSON: Kembalikan MURNI JSON dalam blok:
%sjson
{
  "judul_presentasi": "Judul",
  "rekomendasi_tema": "bisnis",
  "slides": [{"slide_type": "title", "title": "Judul", "content": "Sub-judul"}]
}
%s
- DOKUMEN: Kembalikan konten dalam blok:
%sdocument
# Judul
Isi paragraf...
%s""" % (B3, B3, B3, B3)

# ==========================================
# 2. MANAJER DATABASE
# ==========================================
class DatabaseManager:
    @staticmethod
    @contextmanager
    def get_connection():
        conn = sqlite3.connect(DB_NAME)
        try: yield conn
        finally: conn.close()

    @classmethod
    def init_db(cls):
        with cls.get_connection() as conn:
            c = conn.cursor()
            c.execute('''CREATE TABLE IF NOT EXISTS users (username TEXT PRIMARY KEY, password TEXT)''')
            c.execute('''CREATE TABLE IF NOT EXISTS sessions (session_id TEXT PRIMARY KEY, username TEXT, title TEXT, updated_at TIMESTAMP)''')
            c.execute('''CREATE TABLE IF NOT EXISTS messages (id INTEGER PRIMARY KEY AUTOINCREMENT, session_id TEXT, role TEXT, content TEXT)''')
            conn.commit()

    @staticmethod
    def hash_password(password: str) -> str:
        return hashlib.sha256(password.encode()).hexdigest()

    @classmethod
    def register_user(cls, username: str, password: str) -> bool:
        with cls.get_connection() as conn:
            c = conn.cursor()
            try:
                c.execute("INSERT INTO users (username, password) VALUES (?, ?)", (username, cls.hash_password(password)))
                conn.commit()
                return True
            except sqlite3.IntegrityError:
                return False

    @classmethod
    def authenticate_user(cls, username: str, password: str) -> bool:
        with cls.get_connection() as conn:
            c = conn.cursor()
            c.execute("SELECT password FROM users WHERE username=?", (username,))
            row = c.fetchone()
            return bool(row and row[0] == cls.hash_password(password))

    @classmethod
    def get_user_sessions(cls, username: str) -> List[tuple]:
        with cls.get_connection() as conn:
            c = conn.cursor()
            c.execute("SELECT session_id, title FROM sessions WHERE username=? ORDER BY updated_at DESC", (username,))
            return c.fetchall()

    @classmethod
    def load_session_messages(cls, session_id: str) -> List[Dict[str, Any]]:
        with cls.get_connection() as conn:
            c = conn.cursor()
            c.execute("SELECT role, content FROM messages WHERE session_id=? ORDER BY id ASC", (session_id,))
            rows = c.fetchall()

        msgs = [{"role": "system", "content": SYSTEM_PROMPT}]
        for r, content_str in rows:
            try: msgs.append({"role": r, "content": json.loads(content_str)})
            except: msgs.append({"role": r, "content": content_str})
        return msgs

    @classmethod
    def save_session(cls, session_id: str, username: str, title: str, messages: List[Dict[str, Any]]):
        with cls.get_connection() as conn:
            c = conn.cursor()
            c.execute("INSERT OR REPLACE INTO sessions (session_id, username, title, updated_at) VALUES (?, ?, ?, ?)",
                      (session_id, username, title, datetime.datetime.now()))
            c.execute("DELETE FROM messages WHERE session_id=?", (session_id,))

            for msg in messages:
                if msg["role"] != "system":
                    content = json.dumps(msg["content"])
                    c.execute("INSERT INTO messages (session_id, role, content) VALUES (?, ?, ?)",
                              (session_id, msg["role"], content))
            conn.commit()

    @classmethod
    def delete_session(cls, session_id: str):
        with cls.get_connection() as conn:
            c = conn.cursor()
            c.execute("DELETE FROM sessions WHERE session_id=?", (session_id,))
            c.execute("DELETE FROM messages WHERE session_id=?", (session_id,))
            conn.commit()

@st.cache_resource(show_spinner=False)
def setup_database():
    DatabaseManager.init_db()
    return True

# ==========================================
# 3. UTILITIES & IMPLEMENTASI ALAT (TOOLS)
# ==========================================
class AgentTools:

    @staticmethod
    def _parse_ddg_html(soup) -> List[tuple]:
        hasil = []
        for res in soup.select("div.result")[:6]:
            a = res.select_one("a.result__a")
            sn = res.select_one(".result__snippet")
            if not a and not sn: continue
            url = a.get("href", "") if a else ""
            if url.startswith("//"): url = "https:" + url
            if "uddg=" in url:
                try: url = unquote(parse_qs(urlparse(url).query)["uddg"][0])
                except: pass
            judul = a.get_text(" ", strip=True) if a else ""
            snippet = sn.get_text(" ", strip=True) if sn else ""
            if judul or snippet: hasil.append((judul, url, snippet))
        return hasil

    @staticmethod
    def _parse_bing(soup) -> List[tuple]:
        hasil = []
        for li in soup.select("li.b_algo")[:6]:
            a = li.select_one("h2 a")
            p = li.select_one("p")
            if not a: continue
            judul = a.get_text(" ", strip=True)
            url = a.get("href", "")
            snippet = p.get_text(" ", strip=True) if p else ""
            if judul and url.startswith("http"): hasil.append((judul, url, snippet))
        return hasil

    @staticmethod
    def _fallback_wikipedia(query: str) -> List[tuple]:
        for lang in ("id", "en"):
            try:
                s = HTTP.get(f"https://{lang}.wikipedia.org/w/api.php",
                             params={"action": "opensearch", "search": query, "limit": 1, "format": "json"},
                             timeout=10).json()
                if not s[1]: continue
                title = s[1][0]
                summ = HTTP.get(
                    f"https://{lang}.wikipedia.org/api/rest_v1/page/summary/{quote(title.replace(' ', '_'))}",
                    timeout=10).json()
                extract = summ.get("extract", "")
                url = summ.get("content_urls", {}).get("desktop", {}).get("page", "")
                if extract: return [(f"Wikipedia: {title}", url, extract)]
            except: continue
        return []

    @staticmethod
    def cari_informasi_web(query: str) -> str:
        hasil = []
        try:
            r = HTTP.get("https://html.duckduckgo.com/html/", params={"q": query}, timeout=15)
            if r.status_code == 200:
                hasil = AgentTools._parse_ddg_html(BeautifulSoup(r.text, "html.parser"))
        except: pass

        if not hasil:
            try:
                r = HTTP.post("https://lite.duckduckgo.com/lite/", data={"q": query}, timeout=15)
                if r.status_code == 200:
                    soup = BeautifulSoup(r.text, "html.parser")
                    for tr in soup.find_all("tr"):
                        td = tr.find("td", class_="result-snippet")
                        if td: hasil.append(("", "", td.get_text(" ", strip=True)))
                    hasil = hasil[:6]
            except: pass

        if not hasil:
            try:
                r = HTTP.get("https://www.bing.com/search", params={"q": query, "count": 8}, timeout=15)
                if r.status_code == 200:
                    hasil = AgentTools._parse_bing(BeautifulSoup(r.text, "html.parser"))
            except: pass

        if not hasil:
            hasil = AgentTools._fallback_wikipedia(query)

        if not hasil:
            return f"[Log Sistem] Pencarian gagal menemukan info untuk '{query}'. Mohon analisis secara mandiri atau beri tahu pengguna bahwa data spesifik tidak ditemukan."

        out = [f'HASIL PENCARIAN WEB untuk "{query}":']
        for i, (judul, url, snippet) in enumerate(hasil[:6], 1):
            baris = f"{i}. Judul: {judul}" if judul else f"{i}."
            if url: baris += f"\n   URL: {url}"
            if snippet: baris += f"\n   Isi: {snippet}"
            out.append(baris)
        out.append("\n(SISTEM: Rangkum dan susun informasi di atas ke dalam bahasa Anda sendiri. Jangan pernah menampilkan log ini secara mentah kepada pengguna.)")
        return "\n".join(out)

    @staticmethod
    def cari_gambar(query: str) -> str:
        try:
            for lang in ("id", "en"):
                s = HTTP.get(f"https://{lang}.wikipedia.org/w/api.php",
                             params={"action": "opensearch", "search": query, "limit": 1, "format": "json"},
                             timeout=10).json()
                if not s[1]: continue
                title = s[1][0]
                summ = HTTP.get(
                    f"https://{lang}.wikipedia.org/api/rest_v1/page/summary/{quote(title.replace(' ', '_'))}",
                    timeout=10).json()
                img = (summ.get("originalimage") or summ.get("thumbnail") or {}).get("source")
                if img: return f"Pesan Sistem: Foto '{title}' ditemukan. Tampilkan ke pengguna dengan format: ![{title}]({img})"
            return f"Pesan Sistem: Tidak menemukan foto untuk '{query}'."
        except Exception as e:
            return f"Gagal mencari gambar: {str(e)}"

    @staticmethod
    def ambil_transkrip_youtube(video_url: str) -> str:
        try:
            from youtube_transcript_api import YouTubeTranscriptApi
        except ImportError:
            return "Pesan Sistem: Library youtube_transcript_api belum terinstal."
        try:
            m = re.search(r"(?:v=|youtu\.be/|shorts/)([\w-]{11})", video_url)
            if not m: return "Pesan Sistem: URL YouTube tidak valid."
            video_id = m.group(1)
            try:
                ytt = YouTubeTranscriptApi()
                try:
                    transkrip = ytt.fetch(video_id, languages=["id", "en"])
                except:
                    daftar = ytt.list(video_id)
                    if not daftar.transcripts:
                        return "Pesan Sistem: Video tidak memiliki subtitle."
                    transkrip = ytt.fetch(video_id, languages=[daftar.transcripts[0].language_code])
                teks = " ".join([s.text for s in transkrip])
            except (TypeError, AttributeError):
                data = YouTubeTranscriptApi.get_transcript(video_id, languages=["id", "en"])
                teks = " ".join([t["text"] for t in data])
            return f"Transkrip Video YouTube:\n{teks[:12000]}"
        except Exception as e:
            return f"Gagal mengambil transkrip: {str(e)}"

    @staticmethod
    def eksekusi_python(kode: str) -> str:
        import sys
        old_stdout = sys.stdout
        redirected_output = io.StringIO()
        try:
            sys.stdout = redirected_output
            local_scope = {}
            exec(kode, {}, local_scope)
            output = redirected_output.getvalue()
            return f"Hasil Output:\n{output}" if output else f"Eksekusi Sukses. Variabel: {local_scope}"
        except Exception as e:
            return f"Error: {str(e)}"
        finally:
            sys.stdout = old_stdout

    @staticmethod
    def hitung_matematika(ekspresi: str) -> str:
        try:
            ekspresi = ekspresi.replace(",", ".")
            if not re.match(r'^[\d+\-*/().%\s]+$', ekspresi):
                return "Pesan Sistem: Karakter tidak aman."
            if len(ekspresi) > 200:
                return "Pesan Sistem: Terlalu panjang."
            hasil = eval(ekspresi, {"__builtins__": {}}, {})
            return f"Hasil dari {ekspresi} adalah {hasil}"
        except Exception as e:
            return f"Gagal menghitung ({str(e)})."


class MediaUtils:
    @staticmethod
    @st.cache_data(show_spinner=False)
    def konversi_gambar_ke_base64(uploaded_file) -> Optional[str]:
        if uploaded_file is not None: return base64.b64encode(uploaded_file.read()).decode('utf-8')
        return None

    @staticmethod
    @st.cache_data(show_spinner=False)
    def ekstrak_teks_dari_dokumen(uploaded_file) -> str:
        teks_hasil = ""
        nama_file = uploaded_file.name.lower()
        try:
            if nama_file.endswith('.pdf'):
                from pypdf import PdfReader
                reader = PdfReader(uploaded_file)
                for page in reader.pages:
                    teks = page.extract_text()
                    if teks: teks_hasil += teks + "\n"
            elif nama_file.endswith('.txt'):
                teks_hasil = uploaded_file.read().decode("utf-8")
            elif nama_file.endswith('.docx'):
                from docx import Document
                doc = Document(uploaded_file)
                for para in doc.paragraphs: teks_hasil += para.text + "\n"
            return teks_hasil.strip()
        except Exception as e:
            return f"Gagal membaca dokumen: {str(e)}"

    @staticmethod
    def buat_file_word(riwayat_pesan: List[Dict[str, Any]]) -> io.BytesIO:
        from docx import Document
        doc = Document()
        doc.add_heading('Lagøs AI Agent - Analisis Laporan', 0)
        for msg in riwayat_pesan:
            if msg["role"] in ["system", "tool"]: continue
            if msg["role"] == "assistant" and isinstance(msg.get("content"), dict) and "tool_calls" in msg: continue
            role_title = "User" if msg["role"] == "user" else "Lagøs AI"
            doc.add_heading(f"{role_title}", level=2)
            content = msg["content"]
            text_content = next((item["text"] for item in content if item["type"] == "text"), "") if isinstance(content, list) else str(content)
            if not text_content or text_content == "None": continue
            for line in text_content.split('\n'):
                line = line.strip()
                if not line: continue
                if line.startswith('# '): doc.add_heading(line[2:], 3)
                elif line.startswith('- '): doc.add_paragraph(line[2:], style='List Bullet')
                else: doc.add_paragraph(line)
            doc.add_paragraph("\n" + "_" * 40 + "\n")
        bio = io.BytesIO()
        doc.save(bio)
        bio.seek(0)
        return bio

    @staticmethod
    def ekstrak_dokumen(teks: str) -> Optional[str]:
        if not teks: return None
        match = re.search(r'`{3}document\n(.*?)\n`{3}', teks, re.DOTALL | re.IGNORECASE)
        return match.group(1) if match else None

    @staticmethod
    def buat_dokumen_docx(konten: str) -> io.BytesIO:
        from docx import Document
        doc = Document()
        for line in konten.split('\n'):
            line = line.strip()
            if not line: continue
            if line.startswith('# '): doc.add_heading(line[2:], level=1)
            elif line.startswith('## '): doc.add_heading(line[3:], level=2)
            elif line.startswith('### '): doc.add_heading(line[4:], level=3)
            elif line.startswith('- '): doc.add_paragraph(line[2:], style='List Bullet')
            else: doc.add_paragraph(line)
        bio = io.BytesIO()
        doc.save(bio)
        bio.seek(0)
        return bio

    @staticmethod
    def buat_dokumen_pdf(konten: str) -> io.BytesIO:
        try:
            from fpdf import FPDF
        except ImportError:
            raise ImportError("Fitur PDF diblokir karena library fpdf2 belum diinstal.")
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("helvetica", size=12)
        for line in konten.split('\n'):
            line = line.strip()
            if not line:
                pdf.ln(5); continue
            if line.startswith('# '):
                pdf.set_font("helvetica", style="B", size=16)
                pdf.multi_cell(0, 10, text=line[2:])
                pdf.set_font("helvetica", size=12)
            elif line.startswith('## '):
                pdf.set_font("helvetica", style="B", size=14)
                pdf.multi_cell(0, 10, text=line[3:])
                pdf.set_font("helvetica", size=12)
            elif line.startswith('### '):
                pdf.set_font("helvetica", style="B", size=12)
                pdf.multi_cell(0, 8, text=line[4:])
                pdf.set_font("helvetica", size=12)
            elif line.startswith('- '):
                pdf.multi_cell(0, 8, text=f"• {line[2:]}")
            else:
                pdf.multi_cell(0, 8, text=line)
        bio = io.BytesIO(pdf.output())
        return bio

    @staticmethod
    def generate_title_from_messages(messages: List[Dict[str, Any]]) -> str:
        for msg in messages:
            if msg["role"] == "user":
                content = msg["content"]
                text = next((item["text"] for item in content if item["type"] == "text"), "") if isinstance(content, list) else str(content)
                text = text.strip()
                return text[:25] + "..." if len(text) > 25 else (text if text else "Obrolan Baru")
        return "Obrolan Baru"

    @staticmethod
    def _ambil_via_jina(url: str) -> Optional[str]:
        try:
            r = HTTP.get(f"https://r.jina.ai/{url}", timeout=25)
            if r.status_code == 200 and len(r.text) > 200:
                teks = r.text
                for prefix in ["Title:", "URL Source:", "Markdown Content:", "Published Time:"]:
                    idx = teks.find(prefix)
                    if idx != -1:
                        nl = teks.find("\n", idx)
                        if nl != -1: teks = teks[:idx] + teks[nl + 1:]
                teks = re.sub(r'!\[.*?\]\(.*?\)', '', teks)
                teks = re.sub(r'\[([^\]]+)\]\((.*?)\)', r'\1 (\2)', teks)
                return teks[:15000].strip()
        except: pass
        return None

    @staticmethod
    def _ambil_via_allorigins(url: str) -> Optional[str]:
        try:
            r = HTTP.get(f"https://api.allorigins.win/raw?url={quote(url, safe='')}", timeout=20)
            if r.status_code == 200 and len(r.text) > 200: return r.text
        except: pass
        return None

    @staticmethod
    def _ambil_langsung(url: str) -> Optional[str]:
        try:
            r = HTTP.get(url, timeout=15, allow_redirects=True)
            if r.status_code in [403, 401, 406, 429]: return None
            r.raise_for_status()
            return r.text
        except: return None

    @staticmethod
    def _ekstrak_konten_bersih(html: str, url_asal: str) -> str:
        soup = BeautifulSoup(html, 'html.parser')
        for el in soup(['script', 'style', 'nav', 'footer', 'header', 'aside', 'noscript']):
            el.decompose()
        main_el = None
        for sel in ['article', 'main', '[role="main"]', '.post', '.article', '.content']:
            main_el = soup.select_one(sel)
            if main_el and len(main_el.get_text(strip=True)) > 200: break
        if not main_el: main_el = soup.body or soup
        daftar_gambar = []
        for img in main_el.find_all('img'):
            src = img.get('src') or img.get('data-src')
            if not src: continue
            src = urljoin(url_asal, src)
            if any(ext in src.lower() for ext in ['.svg', 'icon', 'logo', 'avatar']): continue
            alt = img.get('alt', '').strip() or "Gambar"
            if len(daftar_gambar) < 10: daftar_gambar.append(f"- ![{alt}]({src})")
        teks_bagian = []
        for el in main_el.find_all(['h1', 'h2', 'h3', 'h4', 'p', 'li', 'blockquote']):
            t = el.get_text(' ', strip=True)
            if not t or len(t) < 3: continue
            if el.name.startswith('h'):
                level = '#' * int(el.name[1])
                teks_bagian.append(f"{level} {t}")
            elif el.name == 'li':
                teks_bagian.append(f"- {t}")
            else:
                teks_bagian.append(t)
        hasil = "\n\n".join(teks_bagian)[:12000].strip()
        if daftar_gambar: hasil += "\n\n[GAMBAR DI HALAMAN:]\n" + "\n".join(daftar_gambar)
        return hasil

    @staticmethod
    def ambil_teks_dari_link(url: str) -> str:
        try:
            if not url.startswith('http'): url = 'https://' + url
            html = MediaUtils._ambil_via_jina(url)
            if html: return f"[ISI HALAMAN: {url}]\n{html[:12000]}"
            html = MediaUtils._ambil_langsung(url)
            if not html: html = MediaUtils._ambil_via_allorigins(url)
            if not html: return f"Pesan Sistem: Gagal akses {url}. Website mungkin diblokir."
            if len(html) < 200: return f"Pesan Sistem: {url} kosong."
            konten = MediaUtils._ekstrak_konten_bersih(html, url)
            if not konten or len(konten) < 100:
                soup = BeautifulSoup(html, 'html.parser')
                for s in soup(['script', 'style']): s.decompose()
                konten = soup.get_text(' | ', strip=True)[:12000]
            if not konten: return f"Pesan Sistem: {url} tidak ada teks (SPA/JS penuh)."
            return f"[ISI HALAMAN: {url}]\n{konten}"
        except Exception as e:
            return f"Error Link {url}: {str(e)}"

    @staticmethod
    def ekstrak_kode_html(teks: str) -> Optional[str]:
        if not teks: return None
        match = re.search(r'`{3}html\n(.*?)\n`{3}', teks, re.DOTALL | re.IGNORECASE)
        return match.group(1) if match else None

    @staticmethod
    def ekstrak_json_ppt(teks: str) -> Optional[dict]:
        if not teks: return None
        match = re.search(r'`{3}json\n(.*?)\n`{3}', teks, re.DOTALL | re.IGNORECASE)
        if match:
            try: return json.loads(match.group(1))
            except: pass
        return None

    @staticmethod
    def buat_file_ppt(data_json: dict) -> io.BytesIO:
        from pptx import Presentation
        tema_pilihan = data_json.get("rekomendasi_tema", "bisnis").lower()
        peta = {"bisnis": "tema_bisnis.pptx", "kreatif": "tema_kreatif.pptx", "akademik": "tema_akademik.pptx", "gelap": "tema_gelap.pptx"}
        file_template = peta.get(tema_pilihan, "tema_bisnis.pptx")
        prs = Presentation(file_template) if os.path.exists(file_template) else Presentation()
        for slide_data in data_json.get("slides", []):
            stype = slide_data.get("slide_type", "content")
            if stype == "title":
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                try:
                    slide.shapes.title.text = slide_data.get("title", "")
                    slide.placeholders[1].text = slide_data.get("content", "")
                except: pass
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                try:
                    slide.shapes.title.text = slide_data.get("title", "")
                    tf = slide.placeholders[1].text_frame
                    content = slide_data.get("content", [])
                    if isinstance(content, list):
                        for i, poin in enumerate(content):
                            if i == 0: tf.text = poin
                            else: tf.add_paragraph().text = poin
                    else: tf.text = str(content)
                except: pass
        bio = io.BytesIO()
        prs.save(bio)
        bio.seek(0)
        return bio


class MarketUtils:
    @staticmethod
    def ambil_data_pasar(simbol_ticker: str) -> str:
        try:
            if not simbol_ticker.endswith("-USD") and len(simbol_ticker) <= 5:
                if simbol_ticker.upper() in ['BTC', 'ETH', 'SOL', 'BNB', 'XRP', 'ADA', 'DOGE']:
                    simbol_ticker = f"{simbol_ticker.upper()}-USD"
            ticker = yf.Ticker(simbol_ticker)
            hist = ticker.history(period="5d")
            if hist.empty: return f"Pesan Sistem: Data '{simbol_ticker}' tidak ditemukan."
            data_str = hist[['Open', 'High', 'Low', 'Close', 'Volume']].to_string()
            return f"Data 5 Hari Terakhir {simbol_ticker}:\n{data_str}"
        except Exception as e: return f"Gagal: {str(e)}"


# ==========================================
# 4. KOMPONEN UI (PROFESIONAL)
# ==========================================
def inject_custom_css():
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@300;400;500;600;700;800&family=Space+Grotesk:wght@500;600;700&display=swap');

    :root{
      --bg:#0b0e14; --surface:#10141f; --card:#141927; --input:#151a28;
      --border:rgba(255,255,255,.07); --text:#e8ecf4; --muted:#8b93a7;
      --acc1:#7c5cff; --acc2:#22d3ee;
      --grad:linear-gradient(135deg,#7c5cff 0%,#22d3ee 100%);
    }

    html,body,[data-testid="stAppViewContainer"],section.main{background:var(--bg) !important;}
    html,body,[class*="css"]{font-family:'Plus Jakarta Sans',sans-serif;}
    h1,h2,h3{font-family:'Space Grotesk','Plus Jakarta Sans',sans-serif;}
    #MainMenu,footer{visibility:hidden;}
    header[data-testid="stHeader"]{background:transparent;}

    ::-webkit-scrollbar{width:8px;height:8px;}
    ::-webkit-scrollbar-thumb{background:#2a3145;border-radius:8px;}
    ::-webkit-scrollbar-track{background:transparent;}

    @keyframes fadeUp{from{opacity:0;transform:translateY(10px);}to{opacity:1;transform:none;}}
    @keyframes pulse{0%,100%{opacity:1;}50%{opacity:.3;}}

    .brand-bar{display:flex;justify-content:space-between;align-items:center;padding:4px 2px 20px;animation:fadeUp .5s ease;}
    .brand-logo{display:flex;align-items:center;gap:10px;font-family:'Space Grotesk';font-weight:700;font-size:1.2rem;color:#fff;}
    .brand-logo .dot{width:36px;height:36px;border-radius:11px;background:var(--grad);display:flex;align-items:center;justify-content:center;font-size:18px;box-shadow:0 4px 20px rgba(124,92,255,.45);}
    .status-pill{display:flex;align-items:center;gap:8px;padding:6px 14px;border-radius:999px;background:rgba(34,197,94,.1);border:1px solid rgba(34,197,94,.3);color:#4ade80;font-size:.75rem;font-weight:700;}
    .status-pill .pulse{width:7px;height:7px;border-radius:50%;background:#4ade80;animation:pulse 1.4s infinite;}

    .hero{text-align:center;padding:30px 0 10px;animation:fadeUp .6s ease;}
    .hero-badge{display:inline-block;padding:6px 16px;border-radius:999px;background:rgba(124,92,255,.12);border:1px solid rgba(124,92,255,.35);color:#b7a6ff;font-size:.75rem;font-weight:700;letter-spacing:.08em;text-transform:uppercase;margin-bottom:16px;}
    .hero-title{font-size:2.6rem;font-weight:800;color:#fff;margin:0;}
    .hero-title span{background:var(--grad);-webkit-background-clip:text;-webkit-text-fill-color:transparent;}
    .hero-sub{color:var(--muted);font-size:.95rem;max-width:440px;margin:10px auto 26px;}

    [data-testid="stSidebar"]{background:var(--surface) !important;border-right:1px solid var(--border);}
    .profile-card{display:flex;gap:12px;align-items:center;padding:14px;border-radius:16px;background:linear-gradient(145deg,rgba(124,92,255,.14),rgba(34,211,238,.07));border:1px solid rgba(124,92,255,.3);margin-bottom:8px;}
    .profile-avatar{width:42px;height:42px;border-radius:50%;display:flex;align-items:center;justify-content:center;background:var(--grad);color:#fff;font-weight:800;font-size:1.1rem;flex-shrink:0;}
    .profile-name{color:#fff;font-weight:700;font-size:.95rem;}
    .profile-role{color:var(--muted);font-size:.72rem;letter-spacing:.04em;}
    .side-label{font-size:.66rem;letter-spacing:.16em;text-transform:uppercase;color:var(--muted);font-weight:800;margin:16px 0 8px;}

    [data-testid="stSidebar"] .stButton>button{background:transparent !important;border:none !important;border-radius:12px !important;color:var(--text) !important;justify-content:flex-start;box-shadow:none !important;}
    [data-testid="stSidebar"] .stButton>button:hover{background:rgba(255,255,255,.06) !important;}
    [data-testid="stSidebar"] .stButton>button[kind="primary"]{background:rgba(124,92,255,.16) !important;color:#b7a6ff !important;font-weight:700 !important;}
    [data-testid="stSidebar"] .stButton>button p{white-space:nowrap !important;overflow:hidden !important;text-overflow:ellipsis !important;text-align:left !important;}

    .stButton>button{background:var(--card);border:1px solid var(--border);color:var(--text);border-radius:12px;font-weight:600;}
    .stButton>button:hover{border-color:rgba(124,92,255,.5);}
    .stButton>button[kind="primary"]{background:var(--grad) !important;border:none !important;color:#fff !important;font-weight:700;box-shadow:0 6px 22px rgba(124,92,255,.35);}
    .stDownloadButton>button{background:var(--grad) !important;border:none !important;color:#fff !important;border-radius:12px;font-weight:700;}

    .stTextInput input,.stTextInput>div>input{background:var(--input) !important;border:1px solid var(--border) !important;border-radius:12px;color:var(--text);}
    .stTabs [data-testid="stTabsSelectionBar"]{background:var(--grad);height:2px;}
    .stTabs button{color:var(--muted);font-weight:700;}
    .stTabs button[aria-selected="true"]{color:#fff;}

    [data-testid="stChatMessage"]{background:var(--card);border:1px solid var(--border);border-radius:18px;padding:16px 18px;margin-bottom:12px;animation:fadeUp .45s ease both;}
    .user-bubble{display:flex;justify-content:flex-end;margin:14px 0;animation:fadeUp .45s ease both;}
    .user-bubble .inner{max-width:82%;padding:12px 20px;border-radius:18px 18px 4px 18px;background:linear-gradient(135deg,#7c5cff,#5a3df0);color:#fff;box-shadow:0 6px 24px rgba(124,92,255,.35);white-space:pre-wrap;line-height:1.5;}

    .agent-chip{display:inline-flex;align-items:center;gap:8px;margin:2px 10px 2px 0;padding:6px 14px;border-radius:999px;background:rgba(34,211,238,.08);border:1px solid rgba(34,211,238,.28);color:#7dd8ea;font-size:.73rem;font-weight:700;animation:fadeUp .4s ease both;}
    .agent-chip .pulse{width:6px;height:6px;border-radius:50%;background:#22d3ee;animation:pulse 1.2s infinite;}

    [data-testid="stChatInput"] textarea{background:var(--input) !important;border:1px solid var(--border) !important;border-radius:999px !important;padding:14px 20px !important;box-shadow:none !important;color:var(--text);}
    [data-testid="stChatInput"] textarea:focus{border-color:transparent !important;box-shadow:0 0 0 2px rgba(124,92,255,.6) !important;}

    pre,.stCodeBlock{background:#0f1320 !important;border:1px solid var(--border);border-radius:12px;}
    .stAlert{border-radius:12px !important;border:1px solid var(--border) !important;}
    </style>
    """, unsafe_allow_html=True)


def render_brand_bar():
    st.markdown("""
    <div class="brand-bar">
      <div class="brand-logo"><span class="dot">🤖</span> Lagøs <span style="color:#8b93a7;font-weight:500;">AI</span></div>
      <div class="status-pill"><span class="pulse"></span> Sistem Aktif</div>
    </div>""", unsafe_allow_html=True)


def render_profile_card(username):
    nama = str(username) if username else "Guest"
    initial = nama[0].upper() if nama else "G"
    st.markdown(f"""
    <div class="profile-card">
      <div class="profile-avatar">{html_escape(initial)}</div>
      <div><div class="profile-name">{html_escape(nama)}</div>
      <div class="profile-role">ANALYST • AKTIF</div></div>
    </div>""", unsafe_allow_html=True)


def render_agent_chip(nama):
    st.markdown(f'<div class="agent-chip"><span class="pulse"></span> ⚙️ {html_escape(nama)}</div>', unsafe_allow_html=True)

# Fungsi pembersih disederhanakan, hanya menghapus token sistem mesin saja (bukan logika merusak yang agresif)
def bersihkan_teks_response(teks):
    if not teks: return ""
    teks = re.sub(r'<\|.*?\|>', '', teks)
    teks = re.sub(r'<atem:.*?(</atem:function_calls>|$)', '', teks, flags=re.DOTALL)
    teks = re.sub(r'<function_calls>.*?(</function_calls>|$)', '', teks, flags=re.DOTALL)
    return teks.strip()

def render_hero_login():
    st.markdown("""
    <div class="hero">
      <div class="hero-badge">✦ Sistem Analitik Otonom</div>
      <h1 class="hero-title">Lagøs <span>AI Agent</span></h1>
      <p class="hero-sub">Riset web multi-sumber, analisis pasar, dan otomatisasi dokumen — dalam satu percakapan.</p>
    </div>""", unsafe_allow_html=True)


def inject_auto_scroll():
    components.html("""
        <script>
            setTimeout(function() {
                var parentDoc = window.parent.document;
                var marker = parentDoc.getElementById('bottom-marker');
                if (marker) marker.scrollIntoView({behavior: 'auto', block: 'end'});
            }, 300);
        </script>
    """, height=0)


@st.dialog("🌐 Web App Preview", width="large")
def render_webapp_modal(html_code: str):
    st.info("💡 Interaksi dengan Web App di bawah ini.")
    injection = "<base target='_blank'><script>document.addEventListener('click', function(e) { var t = e.target.closest('a'); if(t && t.href) { t.setAttribute('target', '_blank'); } });</script>"
    if re.search(r'<head[^>]*>', html_code, re.IGNORECASE):
        html_code = re.sub(r'(<head[^>]*>)', r'\1\n' + injection, html_code, count=1, flags=re.IGNORECASE)
    else: html_code = injection + "\n" + html_code
    components.html(html_code, height=600, scrolling=True)


def init_session_state():
    defaults = {
        "logged_in": False, "username": "", "current_session_id": None,
        "messages": [{"role": "system", "content": SYSTEM_PROMPT}],
        "temp_image": None, "temp_doc": None, "uploader_key": 0,
        "token_usage": 0
    }
    for key, val in defaults.items():
        if key not in st.session_state: st.session_state[key] = val


URL_REGEX = re.compile(
    r'https?://[^\s<>\\"\'\)\]]+'
    r'|'
    r'\b(?:[a-zA-Z0-9-]+\.)+(?:com|net|org|id|co|io|ai|info|gov|edu)\b[^\s<>\\"\'\)\]]*',
    re.IGNORECASE
)

FILE_EXT_SKIP = {'.pdf', '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls', '.txt', '.csv',
                 '.zip', '.rar', '.7z', '.png', '.jpg', '.jpeg', '.gif', '.mp4', '.mp3',
                 '.mov', '.exe', '.apk', '.iso', '.dmg'}

def _apakah_url_valid(url: str) -> bool:
    u = url.lower()
    if any(u.endswith(ext) for ext in FILE_EXT_SKIP): return False
    return '.' in url and len(url) > 7


def main():
    st.set_page_config(page_title="Lagøs AI Agent", page_icon="🤖", layout="centered", initial_sidebar_state="expanded")
    inject_custom_css()
    setup_database()
    init_session_state()

    cookie_manager = stx.CookieManager(key="cookie_manager")
    cookie_logged_in = cookie_manager.get("is_logged_in")
    cookie_username = cookie_manager.get("saved_username")

    if st.session_state.get("del_cookie") == True:
        cookie_manager.delete("is_logged_in", key="del_login_cookie")
        cookie_manager.delete("saved_username", key="del_user_cookie")
        st.session_state.del_cookie = False
        cookie_logged_in = None
        cookie_username = None

    if cookie_logged_in == "True" and not st.session_state.logged_in:
        st.session_state.logged_in = True
        st.session_state.username = cookie_username

    if st.session_state.get("set_cookie") == True:
        expire_date = datetime.datetime.now() + datetime.timedelta(days=7)
        cookie_manager.set("is_logged_in", "True", expires_at=expire_date, key="set_login_cookie")
        cookie_manager.set("saved_username", st.session_state.username, expires_at=expire_date, key="set_user_cookie")
        st.session_state.set_cookie = False

    if not st.session_state.logged_in:
        render_hero_login()
        col1, col2, col3 = st.columns([1, 1.5, 1])
        with col2:
            with st.container(border=True):
                tab_login, tab_register = st.tabs(["🔑 Masuk", "📝 Daftar Baru"])
                with tab_login:
                    log_user = st.text_input("Username", key="log_user")
                    log_pass = st.text_input("Password", type="password", key="log_pass")
                    if st.button("Masuk", use_container_width=True, type="primary"):
                        if DatabaseManager.authenticate_user(log_user, log_pass):
                            st.session_state.logged_in = True
                            st.session_state.username = log_user
                            st.session_state.set_cookie = True
                            st.rerun()
                        else: st.error("Username atau password salah!")
                with tab_register:
                    reg_user = st.text_input("Username Baru", key="reg_user")
                    reg_pass = st.text_input("Password Baru", type="password", key="reg_pass")
                    if st.button("Daftar & Buat Akun", use_container_width=True):
                        if reg_user and reg_pass:
                            if DatabaseManager.register_user(reg_user, reg_pass): st.success("✅ Berhasil mendaftar!")
                            else: st.error("❌ Username sudah dipakai.")
                        else: st.warning("⚠️ Harap isi data!")
        st.stop()

    render_brand_bar()

    with st.sidebar:
        render_profile_card(st.session_state.username)
        st.divider()

        if st.button("➕ Chat Baru", use_container_width=True, type="primary"):
            st.session_state.current_session_id = None
            st.session_state.messages = [{"role": "system", "content": SYSTEM_PROMPT}]
            st.session_state.token_usage = 0
            st.rerun()

        if st.button("🗑️ Hapus Obrolan Aktif", use_container_width=True):
            if st.session_state.current_session_id:
                DatabaseManager.delete_session(st.session_state.current_session_id)
                st.session_state.current_session_id = None
                st.session_state.messages = [{"role": "system", "content": SYSTEM_PROMPT}]
                st.rerun()
            else:
                st.toast("Tidak ada obrolan aktif yang bisa dihapus.")

        st.markdown('<div class="side-label">Riwayat Percakapan</div>', unsafe_allow_html=True)
        sessions = DatabaseManager.get_user_sessions(st.session_state.username)

        if sessions:
            with st.container(height=450, border=False):
                for sess_id, title in sessions:
                    btn_type = "primary" if st.session_state.current_session_id == sess_id else "secondary"
                    if st.button(title, key=f"btn_{sess_id}", use_container_width=True, type=btn_type, help=title):
                        st.session_state.current_session_id = sess_id
                        st.session_state.messages = DatabaseManager.load_session_messages(sess_id)
                        st.session_state.token_usage = 0
                        st.rerun()

        st.divider()
        st.markdown('<div class="side-label">Model AI</div>', unsafe_allow_html=True)
        selected_model = st.selectbox("Pilih model aktif:", list(MODEL_MAPPING.keys()),
                                       format_func=lambda x: MODEL_MAPPING[x], label_visibility="collapsed")

        st.divider()
        st.markdown('<div class="side-label">Statistik Sesi</div>', unsafe_allow_html=True)
        st.info(f"🪙 Est. Token Dipakai: **{st.session_state.token_usage:,}**")

        if len(st.session_state.messages) > 1:
            st.download_button("📥 Unduh Laporan Chat", data=MediaUtils.buat_file_word(st.session_state.messages),
                                file_name="Lagøs_AI_Chat.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                use_container_width=True)

        st.divider()
        if st.button("🚪 Keluar (Logout)", use_container_width=True):
            st.session_state.logged_in = False
            st.session_state.username = ""
            st.session_state.current_session_id = None
            st.session_state.messages = [{"role": "system", "content": SYSTEM_PROMPT}]
            st.session_state.del_cookie = True
            st.rerun()

    # ========== TAMPILAN PESAN ==========
    for idx, message in enumerate(st.session_state.messages):
        if message["role"] in ["system", "tool"]: continue

        if message["role"] == "assistant" and message.get("tool_calls"):
            for t_call in message["tool_calls"]:
                render_agent_chip(t_call.get("function", {}).get("name", "alat"))
            continue

        content = message.get("content", "")
        if not content: continue

        text_disp = next((item["text"] for item in content if item["type"] == "text"), "") if isinstance(content, list) else str(content)

        if message["role"] == "user":
            st.markdown(f'<div class="user-bubble"><div class="inner">{html_escape(text_disp)}</div></div>', unsafe_allow_html=True)
            continue

        with st.chat_message("assistant"):
            st.markdown(text_disp)

            if message["role"] == "assistant":
                html_code = MediaUtils.ekstrak_kode_html(text_disp)
                if html_code:
                    st.write("")
                    if st.button("🚀 Tampilkan Web App", key=f"btn_webapp_{idx}", use_container_width=True):
                        render_webapp_modal(html_code)

                json_ppt = MediaUtils.ekstrak_json_ppt(text_disp)
                if json_ppt:
                    st.write("")
                    ppt_file = MediaUtils.buat_file_ppt(json_ppt)
                    st.download_button("📊 Unduh Presentasi (.PPTX)", data=ppt_file,
                                        file_name=f"{json_ppt.get('judul_presentasi', 'Presentasi_Lagos')}.pptx",
                                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                        key=f"btn_ppt_{idx}", use_container_width=True, type="primary")

                dokumen_teks = MediaUtils.ekstrak_dokumen(text_disp)
                if dokumen_teks:
                    st.write("")
                    col_doc1, col_doc2 = st.columns(2)
                    with col_doc1:
                        docx_file = MediaUtils.buat_dokumen_docx(dokumen_teks)
                        st.download_button(label="📄 Unduh (.DOCX)", data=docx_file, file_name="Dokumen_Lagos.docx",
                                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                            key=f"btn_docx_{idx}", use_container_width=True, type="primary")
                    with col_doc2:
                        try:
                            pdf_file = MediaUtils.buat_dokumen_pdf(dokumen_teks)
                            st.download_button(label="📕 Unduh (.PDF)", data=pdf_file, file_name="Dokumen_Lagos.pdf",
                                                mime="application/pdf", key=f"btn_pdf_{idx}",
                                                use_container_width=True, type="primary")
                        except Exception as e:
                            st.error(str(e))

    st.markdown("<div style='height: 40px'></div>", unsafe_allow_html=True)
    st.markdown("<div id='bottom-marker'></div>", unsafe_allow_html=True)
    inject_auto_scroll()

    # ========== INPUT PESAN ==========
    with st.container():
        uploader_idx = st.session_state.uploader_key
        if st.session_state.get(f"img_{uploader_idx}"):
            st.markdown(f"<div class='file-pill'>📷 Gambar dilampirkan</div>", unsafe_allow_html=True)
        if st.session_state.get(f"doc_{uploader_idx}"):
            st.markdown(f"<div class='file-pill'>📄 Dokumen dilampirkan</div>", unsafe_allow_html=True)

        col_attach, col_input, col_mic = st.columns([1, 8, 1])
        with col_attach:
            with st.popover("➕"):
                st.session_state.temp_image = st.file_uploader("Upload Image", type=["jpg", "png", "jpeg"],
                                                                label_visibility="collapsed", key=f"img_{uploader_idx}")
                st.session_state.temp_doc = st.file_uploader("Upload Doc", type=["pdf", "txt", "docx"],
                                                              label_visibility="collapsed", key=f"doc_{uploader_idx}")
        with col_input:
            prompt_text = st.chat_input("Tanyakan sesuatu...")
        with col_mic:
            audio_bytes = audio_recorder(text="", recording_color="#ff4b4b", neutral_color="#888888",
                                          icon_name="microphone", icon_size="1.8x", key=f"mic_{uploader_idx}")

    prompt = prompt_text
    if audio_bytes and not prompt_text:
        with st.spinner("Menerjemahkan suara..."):
            try:
                prompt = sr.Recognizer().recognize_google(
                    sr.Recognizer().record(sr.AudioFile(io.BytesIO(audio_bytes))), language="id-ID")
            except:
                st.warning("Suara tidak terdengar jelas.")

    if prompt:
        client = OpenAI(base_url=BASE_URL, api_key=API_KEY)

        with st.chat_message("user"): st.markdown(prompt)
        teks_tambahan = ""

        if st.session_state.temp_doc:
            teks_dok = MediaUtils.ekstrak_teks_dari_dokumen(st.session_state.temp_doc)
            if teks_dok: teks_tambahan += f"\n[KONTEN DOKUMEN: {st.session_state.temp_doc.name}]\n{teks_dok}\n[AKHIR DOKUMEN]\n"

        semua_url = []
        for m in URL_REGEX.finditer(prompt):
            u = m.group(0).rstrip('.,;:!?)')
            if _apakah_url_valid(u) and u not in semua_url:
                semua_url.append(u)

        for url in semua_url[:3]:
            teks_tambahan += f"\n[ISI WEBSITE TERKONEKSI: {url}]\n{MediaUtils.ambil_teks_dari_link(url)}\n"

        if st.session_state.temp_image:
            base64_img = MediaUtils.konversi_gambar_ke_base64(st.session_state.temp_image)
            st.session_state.messages.append({
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_img}"}}
                ]
            })
        else:
            st.session_state.messages.append({"role": "user", "content": prompt})

        payload_khusus_api = copy.deepcopy(st.session_state.messages)
        final_prompt_api = f"{teks_tambahan}\n\nPertanyaan/Instruksi Pengguna:\n{prompt}" if teks_tambahan else prompt

        if isinstance(payload_khusus_api[-1]["content"], list):
            payload_khusus_api[-1]["content"][0]["text"] = final_prompt_api
        else:
            payload_khusus_api[-1]["content"] = final_prompt_api

        # ========== AGENT LOOP (DIBATASI ANTI-SPAM) ==========
        MAX_AGENT_LOOPS = 2
        hasil_tool_semua = []
        query_sudah_dicari = set()

        for loop_idx in range(MAX_AGENT_LOOPS):
            try:
                agent_response = panggil_api_dengan_retry(
                    client,
                    model=selected_model,
                    messages=payload_khusus_api,
                    tools=LAGOS_TOOLS,
                    tool_choice="auto",
                    max_tokens=1500
                )

                response_message = agent_response.choices[0].message

                if not response_message.tool_calls:
                    break

                tc_list = response_message.tool_calls[:2]

                msg_assistant_tool = {
                    "role": "assistant",
                    "content": None,
                    "tool_calls": [
                        {
                            "id": tc.id,
                            "type": tc.type,
                            "function": {"name": tc.function.name, "arguments": tc.function.arguments}
                        } for tc in tc_list
                    ]
                }
                
                st.session_state.messages.append(msg_assistant_tool)
                payload_khusus_api.append(msg_assistant_tool)

                for t_call in tc_list:
                    render_agent_chip(t_call.function.name)

                for tool_call in tc_list:
                    func_name = tool_call.function.name
                    try: func_args = json.loads(tool_call.function.arguments)
                    except: func_args = {}

                    hasil_fungsi = "Error: Alat tidak dikenali."

                    if func_name == "ambil_data_pasar":
                        st.info(f"📈 Menganalisis pasar untuk {func_args.get('simbol_ticker', '')}...")
                        hasil_fungsi = MarketUtils.ambil_data_pasar(func_args.get("simbol_ticker", ""))
                    elif func_name == "cari_informasi_web":
                        query_asli = func_args.get("query", "")
                        query_key = query_asli.strip().lower()
                        if query_key in query_sudah_dicari:
                            hasil_fungsi = "Pesan Sistem: Query ini SUDAH dicari. JANGAN ulangi pencarian yang sama."
                        else:
                            query_sudah_dicari.add(query_key)
                            st.info(f"🔍 Mencari: '{query_asli}'...")
                            hasil_fungsi = AgentTools.cari_informasi_web(query_asli)
                    elif func_name == "baca_isi_website":
                        url = func_args.get("url", "")
                        st.info(f"🌐 Membaca situs: {url}...")
                        hasil_fungsi = MediaUtils.ambil_teks_dari_link(url)
                    elif func_name == "cari_gambar":
                        st.info(f"🖼️ Mencari foto: '{func_args.get('query', '')}'...")
                        hasil_fungsi = AgentTools.cari_gambar(func_args.get("query", ""))
                    elif func_name == "ambil_transkrip_youtube":
                        yt_url = func_args.get("video_url", "")
                        st.info(f"🎬 Ekstrak transkrip: {yt_url}...")
                        hasil_fungsi = AgentTools.ambil_transkrip_youtube(yt_url)
                    elif func_name == "eksekusi_python":
                        st.info(f"🐍 Menjalankan skrip Python...")
                        hasil_fungsi = AgentTools.eksekusi_python(func_args.get("kode", ""))
                    elif func_name == "hitung_matematika":
                        st.info(f"🧮 Menghitung: {func_args.get('ekspresi', '')}...")
                        hasil_fungsi = AgentTools.hitung_matematika(func_args.get("ekspresi", ""))

                    hasil_tool_semua.append(str(hasil_fungsi))

                    tool_msg = {
                        "tool_call_id": tool_call.id,
                        "role": "tool",
                        "name": func_name,
                        "content": str(hasil_fungsi),
                    }
                    st.session_state.messages.append(tool_msg)
                    payload_khusus_api.append(tool_msg)

                time.sleep(1.5)

            except Exception as e:
                st.error(f"Error pada loop agent: {str(e)}")
                break

        # ========== STREAMING JAWABAN AKHIR (Penyintesisan Data) ==========
        if st.session_state.messages and st.session_state.messages[-1]["role"] == "assistant" \
                and "tool_calls" not in st.session_state.messages[-1] and not st.session_state.messages[-1].get("content"):
            st.session_state.messages.pop()
            payload_khusus_api.pop()

        with st.chat_message("assistant"):
            placeholder = st.empty()
            full_response = ""

            try:
                # Perhatikan: parameter tools dihilangkan di sini agar model dipaksa murni mengeluarkan Teks Balasan (Synthesize)
                response_stream = panggil_api_dengan_retry(
                    client,
                    model=selected_model,
                    messages=payload_khusus_api,
                    temperature=0.7,
                    max_tokens=4000,
                    stream=True
                )

                for chunk in response_stream:
                    if chunk.choices and len(chunk.choices) > 0:
                        delta = chunk.choices[0].delta.content
                        if delta:
                            full_response += delta
                            placeholder.markdown(bersihkan_teks_response(full_response) + "▌")

                full_response = bersihkan_teks_response(full_response)

                # Pencegahan khusus jika setelah dibersihkan teksnya kosong
                if not full_response.strip():
                    if hasil_tool_semua:
                        full_response = "Saya telah menemukan informasinya dari internet, namun mengalami sedikit kendala saat mencoba merangkum kalimatnya. Silakan coba tanyakan kembali secara lebih spesifik."
                    else:
                        full_response = "Maaf, API model sedang tidak stabil untuk menjawab saat ini. Mohon ulangi lagi pertanyaan Anda. 🙏"

                placeholder.markdown(full_response)
                st.session_state.messages.append({"role": "assistant", "content": full_response})
                st.session_state.token_usage += (len(str(st.session_state.messages)) // 4)

                if st.session_state.current_session_id is None:
                    st.session_state.current_session_id = str(uuid.uuid4())

                DatabaseManager.save_session(
                    st.session_state.current_session_id,
                    st.session_state.username,
                    MediaUtils.generate_title_from_messages(st.session_state.messages),
                    st.session_state.messages
                )

                st.session_state.temp_image = None
                st.session_state.temp_doc = None
                st.session_state.uploader_key += 1

                st.rerun()

            except Exception as e:
                error_msg = str(e)
                if "429" in error_msg:
                    st.error("⏳ API sedang mencapai batas maksimal dari NVIDIA. Silakan tunggu beberapa saat lagi.")
                elif "404" in error_msg:
                    st.error("❌ Kesalahan 404: Model AI sedang tidak tersedia dari server.")
                else:
                    st.error(f"Kesalahan teknis: {error_msg}")

                if st.session_state.messages and st.session_state.messages[-1]["role"] == "user":
                    st.session_state.messages.pop()

if __name__ == "__main__":
    main()
